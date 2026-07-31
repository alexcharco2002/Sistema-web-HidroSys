from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("output")
OUT.mkdir(exist_ok=True)
PPTX_PATH = OUT / "presentacion_hidrosys_capitulos_4_5.pptx"

W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(12, 45, 64)
BLUE = RGBColor(0, 118, 182)
TEAL = RGBColor(0, 150, 136)
MINT = RGBColor(219, 245, 239)
SKY = RGBColor(225, 244, 252)
GREEN = RGBColor(45, 170, 95)
ORANGE = RGBColor(239, 148, 55)
INK = RGBColor(31, 41, 55)
MUTED = RGBColor(86, 101, 115)
LIGHT = RGBColor(245, 248, 250)
WHITE = RGBColor(255, 255, 255)


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def add_text(slide, text, x, y, w, h, size=24, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = "Aptos Display" if size >= 30 else "Aptos"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.72, 0.38, 8.9, 0.72, 30, NAVY, True)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.18), Inches(1.28), Inches(0.06)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = TEAL
    slide.shapes[-1].line.fill.background()
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.35, 10.2, 0.45, 14, MUTED)


def footer(slide, n):
    add_text(slide, "HidroSys | Capítulos IV y V", 0.72, 7.08, 3.6, 0.2, 8.5, MUTED)
    add_text(slide, f"{n:02d}", 12.0, 7.02, 0.55, 0.26, 9, MUTED, align=PP_ALIGN.RIGHT)


def bullet_list(slide, items, x, y, w, h, size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def pill(slide, label, x, y, w, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.44))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Aptos"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return shape


def metric(slide, value, label, x, y, w, accent=TEAL):
    add_text(slide, value, x, y, w, 0.58, 34, accent, True, PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.65, w, 0.5, 12, MUTED, align=PP_ALIGN.CENTER)


def add_band(slide, color, x=0, y=0, w=13.333, h=7.5):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def add_chart(slide, categories, values, x, y, w, h, title=None):
    data = ChartData()
    data.categories = categories
    data.add_series("Cumplimiento", values)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.number_format = '0.00"%"'
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(9)
    val_axis = chart.value_axis
    val_axis.minimum_scale = 0
    val_axis.maximum_scale = 100
    val_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(220, 230, 235)
    return chart_shape


def create():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1
    s = blank(prs)
    add_band(s, NAVY)
    add_band(s, RGBColor(4, 81, 101), 9.45, 0, 3.9, 7.5)
    for i, (cx, cy, r) in enumerate([(10.0, 1.1, 0.9), (11.15, 2.2, 0.45), (12.25, 4.3, 0.7), (10.6, 5.7, 0.36)]):
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(r), Inches(r))
        circ.fill.solid()
        circ.fill.fore_color.rgb = RGBColor(63, 205, 183) if i % 2 else RGBColor(99, 179, 237)
        circ.fill.transparency = 20
        circ.line.fill.background()
    add_text(s, "HidroSys", 0.76, 0.62, 2.4, 0.45, 20, MINT, True)
    add_text(s, "Sistema web para facturación de agua potable con geolocalización de medidores", 0.76, 1.55, 7.55, 1.65, 36, WHITE, True)
    add_text(s, "Comunidad Sanjapamba | Capítulos IV y V", 0.80, 3.46, 6.6, 0.42, 17, RGBColor(190, 231, 230))
    add_text(s, "Jeferson Charco", 0.80, 6.42, 3.3, 0.35, 15, WHITE, True)
    add_text(s, "Defensa de trabajo de titulación", 0.80, 6.78, 3.8, 0.25, 10.5, RGBColor(190, 231, 230))

    # 2
    s = blank(prs)
    add_title(s, "Problema y punto de partida", "El proceso anterior limitaba el control operativo y el acceso oportuno a la información.")
    bullet_list(s, [
        "Facturación gestionada en un sistema de escritorio y procesos manuales.",
        "Dificultad para consultar historial de pagos y estados de cuenta.",
        "Lecturas de medidores dependientes del registro manual en campo.",
        "Ubicación de medidores poco eficiente para planificar recorridos.",
    ], 0.9, 2.0, 5.4, 3.6, 19)
    add_band(s, SKY, 7.1, 1.82, 4.9, 3.9)
    add_text(s, "Necesidad central", 7.55, 2.18, 3.8, 0.35, 18, NAVY, True)
    add_text(s, "Digitalizar la gestión de facturación y facilitar la ubicación de medidores para mejorar lecturas, pagos y atención a afiliados.", 7.55, 2.75, 3.95, 1.7, 26, NAVY, True)
    footer(s, 2)

    # 3
    s = blank(prs)
    add_title(s, "Objetivos del proyecto")
    add_text(s, "Objetivo general", 0.9, 1.75, 3.1, 0.35, 18, TEAL, True)
    add_text(s, "Desarrollar un sistema web para la facturación del servicio de agua potable en la comunidad Sanjapamba, incorporando geolocalización de medidores.", 0.9, 2.22, 5.7, 1.6, 24, NAVY, True)
    add_text(s, "Objetivos específicos", 7.0, 1.75, 3.5, 0.35, 18, TEAL, True)
    bullet_list(s, [
        "Estudiar procesos de facturación para definir requerimientos.",
        "Desarrollar módulos del sistema, incluido Google Maps API.",
        "Evaluar adecuación funcional conforme a ISO/IEC 25010.",
    ], 7.0, 2.28, 5.25, 2.6, 18)
    footer(s, 3)

    # 4
    s = blank(prs)
    add_title(s, "Del proceso manual a los requisitos", "El levantamiento permitió convertir problemas operativos en funcionalidades evaluables.")
    steps = [("Observación", "Proceso real en la JAAP"), ("BPMN", "Modelado de actividades"), ("35 RF", "Requisitos funcionales"), ("12 módulos", "Organización del sistema")]
    for i, (head, body) in enumerate(steps):
        x = 0.95 + i * 3.05
        pill(s, str(i + 1), x, 2.4, 0.58, TEAL if i < 3 else BLUE)
        add_text(s, head, x, 3.0, 2.35, 0.35, 22, NAVY, True)
        add_text(s, body, x, 3.46, 2.3, 0.7, 15, MUTED)
        if i < 3:
            add_text(s, ">", x + 2.38, 3.05, 0.38, 0.3, 18, TEAL, True, PP_ALIGN.CENTER)
    footer(s, 4)

    # 5
    s = blank(prs)
    add_title(s, "Desarrollo de HidroSys", "SCRUM permitió controlar el avance técnico y cerrar el desarrollo con baja desviación.")
    metric(s, "7", "sprints", 1.0, 2.12, 2.1, TEAL)
    metric(s, "218", "horas de esfuerzo", 3.55, 2.12, 2.4, BLUE)
    metric(s, "3,3 %", "desviación del plan", 6.35, 2.12, 2.4, ORANGE)
    add_text(s, "Arquitectura implementada", 1.0, 4.25, 3.8, 0.35, 20, NAVY, True)
    for i, item in enumerate(["React", "FastAPI", "PostgreSQL"]):
        pill(s, item, 1.0 + i * 2.05, 4.88, 1.55, [BLUE, TEAL, GREEN][i])
    add_text(s, "Frontend web, backend API y base de datos relacional para centralizar usuarios, lecturas, facturación y pagos.", 7.25, 4.45, 4.7, 1.0, 19, INK)
    footer(s, 5)

    # 6
    s = blank(prs)
    add_title(s, "Módulos implementados", "El sistema integra la operación administrativa, la lectura de consumos y la atención del afiliado.")
    modules = ["Autenticación", "Usuarios", "Medidores", "Sectores", "Tarifas", "Servicios", "Lecturas", "Facturación", "Pagos", "Multas", "Reportes", "Geolocalización"]
    for i, m in enumerate(modules):
        row, col = divmod(i, 4)
        pill(s, m, 0.9 + col * 3.05, 2.0 + row * 0.85, 2.35, TEAL if i in [2, 6, 11] else BLUE)
    add_text(s, "Roles evaluados", 0.95, 5.35, 2.4, 0.3, 17, NAVY, True)
    add_text(s, "Administrador  |  Cajero  |  Lector  |  Afiliado", 2.85, 5.35, 7.2, 0.3, 17, MUTED)
    footer(s, 6)

    # 7
    s = blank(prs)
    add_title(s, "Geolocalización de medidores", "El módulo conecta la facturación con el territorio real de la comunidad.")
    add_band(s, MINT, 0.9, 1.82, 5.3, 4.5)
    for x, y, color in [(1.55, 2.45, BLUE), (3.7, 2.95, TEAL), (2.55, 4.35, GREEN), (4.9, 4.75, ORANGE)]:
        pin = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.34), Inches(0.34))
        pin.fill.solid()
        pin.fill.fore_color.rgb = color
        pin.line.fill.background()
    add_text(s, "Google Maps API", 1.35, 5.55, 3.4, 0.35, 20, NAVY, True)
    bullet_list(s, [
        "Visualización interactiva de medidores georreferenciados.",
        "Ubicación más rápida durante la lectura de consumo.",
        "Mejor base para planificar recorridos por sector.",
    ], 7.0, 2.2, 5.0, 2.7, 20)
    footer(s, 7)

    # 8
    s = blank(prs)
    add_title(s, "Evaluación de adecuación funcional", "La evaluación se alineó con ISO/IEC 25010 y métricas ISO/IEC 25023.")
    metric(s, "23", "participantes", 1.0, 2.0, 2.4, TEAL)
    metric(s, "3", "miembros administrativos", 3.75, 2.0, 2.5, BLUE)
    metric(s, "20", "afiliados", 6.75, 2.0, 2.2, GREEN)
    bullet_list(s, [
        "Evaluación por roles y módulos funcionales.",
        "Lista de verificación como instrumento de recolección.",
        "Subcaracterísticas: completitud, corrección y pertinencia funcional.",
    ], 1.05, 4.22, 10.5, 1.7, 18)
    footer(s, 8)

    # 9
    s = blank(prs)
    add_title(s, "Resultados por subcaracterística", "Las tres mediciones superan el 93 % y sostienen un alto cumplimiento funcional.")
    add_chart(s, ["Completitud", "Corrección", "Pertinencia"], [97.92, 93.75, 93.75], 1.0, 1.85, 6.7, 4.35)
    add_text(s, "Lectura clave", 8.2, 2.08, 2.8, 0.35, 18, TEAL, True)
    bullet_list(s, [
        "Once de doce módulos alcanzaron 100 % en completitud.",
        "Facturación obtuvo 75 % por una funcionalidad ausente.",
        "OU-01, OU-07 y OU-08 presentaron oportunidades de mejora.",
    ], 8.2, 2.65, 4.0, 2.3, 17)
    footer(s, 9)

    # 10
    s = blank(prs)
    add_title(s, "Resultado global de calidad", "La ponderación de las subcaracterísticas ubicó al sistema en nivel 4 de calidad.")
    add_text(s, "95,14 %", 1.0, 1.88, 4.4, 1.2, 60, TEAL, True)
    add_text(s, "Adecuación funcional global de HidroSys", 1.08, 3.16, 4.6, 0.5, 21, NAVY, True)
    bullet_list(s, [
        "Cumple adecuadamente las funcionalidades establecidas.",
        "La evaluación confirma la utilidad operativa del sistema web.",
        "Los resultados dejan mejoras puntuales para futuras iteraciones.",
    ], 6.6, 2.15, 5.25, 2.3, 20)
    footer(s, 10)

    # 11
    s = blank(prs)
    add_title(s, "Conclusiones principales", "El proyecto resolvió necesidades funcionales concretas de la gestión de agua potable.")
    bullet_list(s, [
        "El análisis del proceso permitió detectar fallas en pagos, lecturas y acceso a información.",
        "HidroSys integró módulos administrativos, operativos y de consulta para los afiliados.",
        "La geolocalización mejora la ubicación de medidores y la planificación de lecturas.",
        "La evaluación ISO/IEC 25010 confirmó un alto grado de cumplimiento funcional.",
    ], 0.95, 1.95, 11.2, 3.8, 20)
    footer(s, 11)

    # 12
    s = blank(prs)
    add_title(s, "Recomendaciones", "Las mejoras propuestas amplían el alcance operativo y la calidad futura del sistema.")
    recs = [
        ("Rutas óptimas", "Directions API o Routes API para ordenar visitas por cercanía."),
        ("App móvil offline", "Flutter para registrar lecturas en campo y sincronizar al recuperar conexión."),
        ("Notificaciones", "Correo, SMS o WhatsApp para alertas de factura, vencimiento y deuda."),
        ("Más calidad", "Evaluar eficiencia de desempeño, seguridad y protección de datos."),
    ]
    for i, (head, body) in enumerate(recs):
        x = 0.95 + (i % 2) * 5.95
        y = 2.0 + (i // 2) * 1.78
        add_text(s, head, x, y, 3.0, 0.32, 19, TEAL if i % 2 == 0 else BLUE, True)
        add_text(s, body, x, y + 0.44, 4.85, 0.75, 16, INK)
    footer(s, 12)

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    create()
